from typing import Any, List, Dict, Optional
from pptx import Presentation

class PresentationParser:
    def __init__(self, pptx_path: str):
        """
        Initializes the PresentationParser with the path to the PowerPoint file.

        Args:
            pptx_path (str): The path to the PowerPoint file.
        """
        self.pptx_path = pptx_path
        self.presentation_data: Optional[List[Dict[str, Any]]] = None
        self.total_slides = 0  
    
    def parse_presentation(self) -> List[Dict[str, Any]]:
        """
        Parses the content of a PowerPoint presentation.

        Returns:
            List[Dict[str, Any]]: A list of dictionaries containing slide information.
        """
        presentation = Presentation(self.pptx_path)
        self.total_slides = len(presentation.slides)
        return [
            {
                "slide_number": slide_number,
                "title": title if (title := next((shape.text.strip() for shape in slide.shapes if hasattr(shape, 'text') and shape == slide.shapes.title), "")) else "",
                "content": [shape.text.strip() for shape in slide.shapes if hasattr(shape, 'text') and shape != slide.shapes.title]
            }
            for slide_number, slide in enumerate(presentation.slides, start=1)
        ]
    
    def print_presentation_data(self) -> None:
        """
        Prints the parsed data of a presentation.

        Returns:
            None
        """
        if self.presentation_data:
            for slide_data in self.presentation_data:
                print(f"Slide {slide_data['slide_number']}\nTitle: {slide_data['title']}\nContent:\n" + ''.join([f"- {content}\n" for content in slide_data['content']]) + '='*40 + '\n')
    
    def get_slide(self, slide_number: int) -> Optional[Dict[str, Any]]:
        """
        Retrieves the content of a specific slide by its slide number.

        Args:
            slide_number (int): The slide number to retrieve.

        Returns:
            Optional[Dict[str, Any]]: A dictionary containing slide information if found, otherwise None.
        """
        if not self.presentation_data:
            self.presentation_data = self.parse_presentation()

        return next(
            (
                slide
                for slide in self.presentation_data
                if slide['slide_number'] == slide_number
            ),
            None,
        )
    
    def process_presentation(self, flag: bool = False) -> None:
        """
        Processes a presentation by parsing the presentation data and optionally printing it.

        Args:
            flag (bool, optional): Flag to determine whether to print the presentation data. Defaults to False.

        Returns:
            None
        """
        self.presentation_data = self.parse_presentation()
        if flag:
            self.print_presentation_data()

def main() -> None:
    """
    Main function to test the PresentationParser class.

    Returns:
        None
    """
    PPTX_PATH = "/mnt/c/Users/AMIR/ramatgan-window/cop1/ExcellentTeam/python-Home-Work/final-exercise-AmirAmmar123/16-Photometry.pptx"
    parser = PresentationParser(PPTX_PATH)
    parser.process_presentation(True)

    slide_number_to_get = 3
    if specific_slide := parser.get_slide(slide_number_to_get):
        print(f"Retrieved Slide {specific_slide['slide_number']}:\nTitle: {specific_slide['title']}\nContent:\n" + ''.join([f"- {content}\n" for content in specific_slide['content']]))
    else:
        print(f"Slide {slide_number_to_get} not found.")
    
    print(f"Total slides in the presentation: {parser.total_slides}")

if __name__ == "__main__":
    main()
